"""
Otto Python Enhancement Demo
=============================
Creates a hybrid slide: CSS background (PNG from Edge) + python-pptx native chart + editable text.
Proves the Python pipeline works with Otto's CSS engine.
"""
import subprocess, os, json, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = os.path.join(os.path.dirname(__file__), '..', 'otto-hybrid-demo.pptx')
TMP = os.path.join(os.path.dirname(__file__), '..', 'temp-hybrid')
EDGE = 'C:\\Program Files (x86)\\Microsoft\\Edge\\Application\\msedge.exe'

os.makedirs(TMP, exist_ok=True)

# CSS background - mesh gradient + orbs, no text
bg_css = ':root{--bg:#050510;--surface:#0a0a18;--accent:#6366f1;--accent2:#a855f7;--ink:#f0f0f8}'+\
 '*{margin:0;padding:0}body{width:1920px;height:1080px;overflow:hidden;background:var(--bg)}'

bg_html = '<!DOCTYPE html><html><head><meta charset="utf-8"><style>' + bg_css + '\n' + '''
.mesh{{position:absolute;width:900px;height:700px;left:5%;top:-10%;border-radius:50%;background:radial-gradient(ellipse,rgba(99,102,241,0.15),transparent 70%);filter:blur(70px)}}
.mesh2{{position:absolute;width:600px;height:500px;right:0;bottom:-5%;border-radius:50%;background:radial-gradient(ellipse,rgba(168,85,247,0.12),transparent 70%);filter:blur(80px)}}
</style></head><body><div class="mesh"></div><div class="mesh2"></div></body></html>'''

bg_path = os.path.join(TMP, 'bg.html')
png_path = os.path.join(TMP, 'bg.png')
with open(bg_path, 'w', encoding='utf-8') as f: f.write(bg_html)

# Screenshot background
print('[otto] Generating CSS background...')
subprocess.run([
    EDGE, '--headless', '--disable-gpu', '--window-size=1920,1080',
    f'--screenshot={png_path}', f'file://{bg_path.replace(chr(92), "/")}'
], timeout=15000, capture_output=True)
print(f'   ✅ Background PNG: {os.path.getsize(png_path)/1024:.0f} KB\n')

# Build PPTX with python-pptx
print('[python] Python-pptx: adding native elements...')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Background + Native Chart + Editable Text
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
slide1.shapes.add_picture(png_path, 0, 0, prs.slide_width, prs.slide_height)

# Title - editable!
txBox = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Otto 混合渲染：CSS 背景 + 原生可编辑图表"
run.font.size = Pt(40)
run.font.bold = True
run.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF8)

# Native bar chart - completely editable in PowerPoint
chart_data = CategoryChartData()
chart_data.categories = ['文档PPT', '数据分析', '飞书同步', '日历任务', '本地安全']
chart_data.add_series('Otto 完成度 (%)', [95, 90, 85, 80, 100])
chart_data.add_series('人工耗时 (h/天)', [4.5, 3.2, 2.8, 1.5, 0.5])

chart_frame = slide1.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(1), Inches(1.8), Inches(7), Inches(4.5),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = 2  # BOTTOM
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = RGBColor(0x63, 0x66, 0xF1)
chart.series[1].format.fill.solid()
chart.series[1].format.fill.fore_color.rgb = RGBColor(0xA8, 0x55, 0xF7)

# Right side: native shapes with editable text
for idx, (label, value, color) in enumerate([
    ("内置技能", "8 个", "6366f1"),
    ("视觉系统", "20 套", "a855f7"),
    ("布局模板", "9 种", "22c55e"),
    ("CSS 香料", "10 道", "f59e0b"),
]):
    y = 1.8 + idx * 1.35
    shape = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(y), Inches(3.8), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x28)
    shape.line.color.rgb = RGBColor(
        int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
    )
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{value}  {label}"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF8)
    p.alignment = PP_ALIGN.CENTER

# Subtitle - editable
txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "✅ 此图表、文字框、色块在 PowerPoint 中完全可编辑 · 可改数据 · 可翻译 · 可换色"
run2.font.size = Pt(16)
run2.font.color.rgb = RGBColor(0x90, 0x90, 0xB0)

# Slide 2: Complex native table
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
slide2.shapes.add_picture(png_path, 0, 0, prs.slide_width, prs.slide_height)

# Editable title
txBox3 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.8))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = "原生可编辑表格：能力矩阵"
run3.font.size = Pt(40)
run3.font.bold = True
run3.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF8)

# Native table
rows, cols = 9, 4
table_shape = slide2.shapes.add_table(
    rows, cols,
    Inches(0.8), Inches(1.8),
    Inches(11.5), Inches(5.2)
)
table = table_shape.table

headers = ['能力', '描述', '状态', '版本']
data = [
    ['文档生成', '一句话 → PPTX/Word', '✅ 生产就绪', 'v8'],
    ['数据分析', '自然语言操作 Excel', '✅ 生产就绪', 'v2'],
    ['飞书同步', '双向消息·日历·审批', '✅ 生产就绪', 'v3'],
    ['CSV 可视化', '10 道 CSS 香料渲染', '✅ 生产就绪', 'v8'],
    ['增量编辑', '单页改写不复写全 deck', '✅ 生产就绪', 'v6'],
    ['来源追溯', '每数绑定出处 source-map', '✅ 生产就绪', 'v6'],
    ['混合渲染', 'CSS 背景+Python 原生对象', '🆕 Python 增强', 'v8'],
    ['Native Charts', 'python-pptx 原生图表', '🆕 Python 增强', 'v8'],
]

for c_idx, h in enumerate(headers):
    cell = table.cell(0, c_idx)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF8)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3E)

for r_idx, row in enumerate(data):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx + 1, c_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(15)
                color = '22c55e' if '✅' in val else ('f59e0b' if '🆕' in val else 'd0d0e0')
                r.font.color.rgb = RGBColor(int(color[0:2],16), int(color[2:4],16), int(color[4:6],16))
        if r_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x28)

prs.save(OUT)
print(f'   ✅ Hybrid PPTX: {OUT} ({os.path.getsize(OUT)/1024:.0f} KB, 2 slides)')
print(f'   Slide 1: CSS background + native bar chart + 4 native shapes + editable text')
print(f'   Slide 2: CSS background + native table (9 rows, fully editable)')

# Cleanup
import shutil
shutil.rmtree(TMP, ignore_errors=True)
