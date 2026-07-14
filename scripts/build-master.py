"""
Otto PPT Master v9 — Production Pipeline
=========================================
CSS backgrounds (Edge headless) + python-pptx native objects.
One command: python scripts/build-master.py
Output: otto-master-deck.pptx — 8 slides, fully editable.
"""
import subprocess, os, sys, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TMP = os.path.join(ROOT, 'temp-master')
OUT = os.path.join(ROOT, 'otto-master-deck.pptx')
EDGE = 'C:\\Program Files (x86)\\Microsoft\\Edge\\Application\\msedge.exe'

os.makedirs(TMP, exist_ok=True)

# ═══════════════════════════════════════
# PHASE 1: CSS Backgrounds
# ═══════════════════════════════════════

VARS = ':root{--bg:#050510;--accent:#6366f1;--accent2:#a855f7;--ink:#f0f0f8}'+\
 '*{margin:0;padding:0}body{width:1920px;height:1080px;overflow:hidden;background:var(--bg)}'

def mesh_bg(extra=''):
    return '<div style="position:absolute;inset:0">'+\
        '<div style="position:absolute;width:1000px;height:800px;left:8%;top:-15%;border-radius:50%;background:radial-gradient(ellipse,rgba(99,102,241,0.16),transparent 70%);filter:blur(70px)"></div>'+\
        '<div style="position:absolute;width:700px;height:600px;right:-3%;bottom:-8%;border-radius:50%;background:radial-gradient(ellipse,rgba(168,85,247,0.12),transparent 70%);filter:blur(80px)"></div>'+\
        '<div style="position:absolute;width:500px;height:500px;left:55%;top:30%;border-radius:50%;background:radial-gradient(ellipse,rgba(99,102,241,0.08),transparent 70%);filter:blur(50px)"></div>'+\
        extra+'</div>'

def noise_layer():
    return '<div style="position:absolute;inset:0;opacity:0.025;pointer-events:none;z-index:99;background-image:url(&quot;data:image/svg+xml,%3Csvg viewBox=\'0 0 256 256\' xmlns=\'http://www.w3.org/2000/svg\'%3E%3Cfilter id=\'n\'%3E%3CfeTurbulence type=\'fractalNoise\' baseFrequency=\'0.85\' numOctaves=\'4\' stitchTiles=\'stitch\'/%3E%3C/filter%3E%3Crect width=\'100%25\' height=\'100%25\' filter=\'url(%23n)\'/%3E%3C/svg%3E&quot;);background-size:256px"></div>'

def orbs_layer():
    return '''<svg viewBox="0 0 1920 1080" style="position:absolute;inset:0" xmlns="http://www.w3.org/2000/svg">
<defs><radialGradient id="o1"><stop offset="0%" stop-color="#6366f1" stop-opacity="0.15"/><stop offset="100%" stop-color="transparent"/></radialGradient>
<filter id="ob"><feGaussianBlur stdDeviation="90"/></filter></defs>
<circle cx="350" cy="280" r="450" fill="url(#o1)" filter="url(#ob)"/>
<circle cx="1550" cy="680" r="380" fill="url(#o1)" filter="url(#ob)" opacity="0.6"/>
</svg>'''

def grid_burst():
    return '<div style="position:absolute;inset:0;opacity:0.4;background-image:linear-gradient(rgba(99,102,241,0.04) 1px,transparent 1px),linear-gradient(90deg,rgba(99,102,241,0.04) 1px,transparent 1px);background-size:60px 60px;mask-image:radial-gradient(ellipse at 70% 50%,black 30%,transparent 70%)"></div>'

bg_configs = {
    'S1_COVER': mesh_bg() + orbs_layer() + noise_layer(),
    'S2_PROBLEM': mesh_bg() + noise_layer(),
    'S3_SOLUTION': mesh_bg() + noise_layer(),
    'S4_CAPABILITIES': mesh_bg() + orbs_layer() + noise_layer(),
    'S5_TRUST': mesh_bg() + noise_layer(),
    'S6_ARCH': mesh_bg() + grid_burst() + noise_layer(),
    'S7_BIGNUM': mesh_bg() + orbs_layer() + noise_layer(),
    'S8_CLOSE': mesh_bg() + orbs_layer() + noise_layer(),
}

print('[phase 1] CSS backgrounds...')
for sid, bg_html_inner in bg_configs.items():
    html = f'<!DOCTYPE html><html><head><meta charset="utf-8"><style>{VARS}</style></head><body>{bg_html_inner}</body></html>'
    hp = os.path.join(TMP, f'{sid}.html')
    pp = os.path.join(TMP, f'{sid}.png')
    with open(hp, 'w', encoding='utf-8') as f: f.write(html)
    subprocess.run([EDGE, '--headless', '--disable-gpu', '--window-size=1920,1080',
        f'--screenshot={pp}', f'file://{hp.replace(chr(92),"/")}'],
        timeout=15000, capture_output=True)
    print(f'  {sid}.png ({os.path.getsize(pp)//1024} KB)')

# ═══════════════════════════════════════
# PHASE 2: Python-pptx Native Build
# ═══════════════════════════════════════

print('\n[phase 2] Python-pptx native build...')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
INK = RGBColor(0xF0, 0xF0, 0xF8)
MUTED = RGBColor(0x90, 0x90, 0xB0)
ACCENT = RGBColor(0x63, 0x66, 0xF1)
ACCENT2 = RGBColor(0xA8, 0x55, 0xF7)

def add_bg(slide, sid):
    p = os.path.join(TMP, f'{sid}.png')
    if os.path.exists(p):
        slide.shapes.add_picture(p, 0, 0, prs.slide_width, prs.slide_height)

def add_title(slide, text, left=1.2, top=0.6, width=11, height=1, size=44, color=INK, align=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for line in text.split('\n'):
        if p.runs: p = tf.add_paragraph(); p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = True
        run.font.color.rgb = color

def add_body(slide, text, left, top, width, height, size=20, color=MUTED, align=PP_ALIGN.LEFT, line_spacing=30):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = Pt(line_spacing)
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.color.rgb = color

def add_overline(slide, text, left=1.2, top=0.35, size=12, color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = True
    run.font.color.rgb = color
    run.font.caps = True

# --- S1: COVER ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S1_COVER')
add_title(s, 'Otto\nYour AI Coworker', top=3.2, size=100)
add_body(s, '值得信赖的企业级 AI 办公引擎', 2, 5.8, 9, 0.6, size=24, color=MUTED, align=PP_ALIGN.CENTER)
add_body(s, 'INVESTOR DECK · 2026', 4, 6.5, 5, 0.3, size=14, color=RGBColor(0x60, 0x60, 0x80), align=PP_ALIGN.CENTER)
print('  S1_COVER')

# --- S2: PROBLEM ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S2_PROBLEM')
add_overline(s, 'THE PROBLEM')
add_title(s, '每个知识工作者\n每天浪费 2.1 小时\n在重复办公任务上', align=PP_ALIGN.LEFT, size=46)
add_body(s, '文档起草 · PPT 排版 · 表格公式\n会议纪要 · 邮件回复 · 数据整理', 1.2, 4.2, 6, 1.2, size=20)
add_body(s, '来源: McKinsey Global Institute 2025', 1.2, 5.6, 5, 0.3, size=14, color=RGBColor(0x70, 0x70, 0x90))
# Right side: mega number
add_body(s, '2.1h', 8, 2.2, 5, 1.5, size=200, color=ACCENT, align=PP_ALIGN.CENTER)
add_body(s, '每日低价值\n重复劳动', 8.5, 4.8, 4, 0.8, size=18, color=MUTED, align=PP_ALIGN.CENTER)
print('  S2_PROBLEM')

# --- S3: SOLUTION ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S3_SOLUTION')
add_overline(s, 'THE SOLUTION')
add_title(s, '把你的工作工具全部 AI 化', align=PP_ALIGN.LEFT, size=40, top=0.55)
cards = [
    ('📄 文档 & PPT', '一句话 → 发布会级演示\n20套视觉 · 9种布局\nHTML→PNG→PPTX 全自动'),
    ('💬 飞书双向同步', '@Otto 即可办公\n消息·日历·文档\n审批·任务全打通'),
    ('🔒 本地优先·安全', '数据不出本地机器\n来源完全可追溯\n检查点可恢复'),
    ('🚀 纯 Node.js', '零外部依赖\n安装即用·跨平台\nElectron+CLI'),
    ('📊 表格 & 数据', '自然语言操作 Excel\n自动分析·图表生成\n公式推导·数据清洗'),
    ('🎨 CSS 视觉引擎', '10道视觉香料\nMesh·Glass·Orbs\nNoise·Kinetic·3D'),
]
for idx, (title, desc) in enumerate(cards):
    col = idx % 3; row = idx // 3
    x = 0.8 + col * 4.1; y = 1.9 + row * 2.7
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.7), Inches(2.3))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x28)
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5E); shape.line.width = Pt(1)
    add_title(s, title, left=x+0.3, top=y+0.2, width=3.1, height=0.5, size=18, color=INK, align=PP_ALIGN.LEFT)
    add_body(s, desc, x+0.3, y+0.8, 3.1, 1.3, size=13, color=MUTED)
print('  S3_SOLUTION')

# --- S4: CAPABILITIES ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S4_CAPABILITIES')
add_overline(s, 'CORE CAPABILITIES')
add_title(s, '开箱即用的 AI 办公能力', align=PP_ALIGN.LEFT, size=40, top=0.55)
# Native chart
cd = CategoryChartData()
cd.categories = ['文档PPT', '数据分析', '飞书同步', '日历任务', '本地安全']
cd.add_series('Otto 完成度 (%)', [95, 90, 85, 80, 100])
cd.add_series('人工耗时 (h/天)', [4.5, 3.2, 2.8, 1.5, 0.5])
cf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.7), Inches(7.8), Inches(5), cd)
ch = cf.chart; ch.has_legend = True; ch.legend.position = 2
ch.series[0].format.fill.solid(); ch.series[0].format.fill.fore_color.rgb = ACCENT
ch.series[1].format.fill.solid(); ch.series[1].format.fill.fore_color.rgb = ACCENT2
# Right side stat shapes
for idx, (label, val, clr) in enumerate([('Built-in Skills', '8', '6366f1'),('Visual Systems', '20', 'a855f7'),('Layouts', '9', '22c55e'),('CSS Spices', '10', 'f59e0b')]):
    y = 1.7 + idx * 1.3
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(y), Inches(3.5), Inches(1.1))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x28)
    sh.line.color.rgb = RGBColor(int(clr[0:2],16), int(clr[2:4],16), int(clr[4:6],16)); sh.line.width = Pt(1.5)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = f'{val}  {label}'
    run.font.size = Pt(20); run.font.bold = True; run.font.color.rgb = INK
print('  S4_CAPABILITIES')

# --- S5: TRUST ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S5_TRUST')
add_overline(s, 'TRUST & SAFETY')
add_title(s, '值得信赖的 AI 同事', align=PP_ALIGN.LEFT, size=40, top=0.55)
# Native table
rows, cols = 6, 3
ts = s.shapes.add_table(rows+1, cols, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5))
tbl = ts.table
for c_idx, h in enumerate(['质量保证', '描述', '状态']):
    cell = tbl.cell(0, c_idx); cell.text = h
    for pp in cell.text_frame.paragraphs:
        for rr in pp.runs: rr.font.bold = True; rr.font.size = Pt(15); rr.font.color.rgb = INK
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3E)
data = [
    ('稳定 slideId', '每页永不改变的标识符', '✅'),
    ('增量编辑', '改一页只再生一页', '✅'),
    ('检查点恢复', '7级 checkpoint 中断可续', '✅'),
    ('来源追溯', '每数绑定出处 source-map', '✅'),
    ('优雅降级', 'Needs-Manual 不假完成', '✅'),
    ('审计脚本', 'audit-pptx.cjs 自动校验', '✅'),
]
for ri, (a, b, c) in enumerate(data):
    for ci, txt in enumerate([a, b, c]):
        cell = tbl.cell(ri+1, ci); cell.text = txt
        for pp in cell.text_frame.paragraphs:
            for rr in pp.runs: rr.font.size = Pt(14)
            if ci == 2:
                for rr in pp.runs: rr.font.color.rgb = RGBColor(0x22, 0xC5, 0x5E)
            else:
                for rr in pp.runs: rr.font.color.rgb = MUTED
        if ri % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x28)
print('  S5_TRUST')

# --- S6: ARCHITECTURE ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S6_ARCH')
add_overline(s, 'ARCHITECTURE')
add_title(s, '五端统一 · 纯 Node.js 引擎', align=PP_ALIGN.LEFT, size=40, top=0.55)
add_body(s, 'Electron 桌面 + CLI + VSCode 插件\n飞书 Bot · Server API 后端\n三端统一体验，安装即用', 1.2, 1.8, 5.5, 2, size=20)
items = [('C', 'CLI Terminal', '开发者原生 · 管道友好'), ('D', 'Desktop App', 'Electron 原生 · GUI 全功能'),
         ('F', 'Feishu Bot', '消息同步 · 日历 · 审批'), ('V', 'VSCode Plugin', 'IDE 内嵌 · 代码+文档'),
         ('S', 'Server Backend', 'REST API · 飞书网关')]
for idx, (letter, name, desc) in enumerate(items):
    y = 1.7 + idx * 1.05
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(y), Inches(5.2), Inches(0.9))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x28)
    sh.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5E); sh.line.width = Pt(1)
    # Letter badge
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(y+0.15), Inches(0.6), Inches(0.6))
    badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT
    tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = letter; run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    add_title(s, name, left=8.5, top=y+0.08, width=3.5, height=0.35, size=17, color=INK, align=PP_ALIGN.LEFT)
    add_body(s, desc, 8.5, y+0.42, 3.5, 0.3, size=13, color=MUTED)
print('  S6_ARCH')

# --- S7: BIG NUMBER ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S7_BIGNUM')
add_body(s, '1', 3.5, 2.5, 6, 3, size=280, color=ACCENT, align=PP_ALIGN.CENTER)
add_body(s, 'SENTENCE → DELIVERABLE PPTX', 3, 5.8, 7, 0.5, size=20, color=MUTED, align=PP_ALIGN.CENTER)
add_body(s, '从一句话到可交付的演示文稿\n不需要设计师 · 不需要模板', 3.5, 6.3, 6, 1, size=18, color=MUTED, align=PP_ALIGN.CENTER)
print('  S7_BIGNUM')

# --- S8: CLOSE ---
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, 'S8_CLOSE')
add_title(s, '让 Otto 成为\n你团队的第 N+1 个成员', top=2.8, size=60)
add_body(s, '纯 Node.js · 开源 · 值得信赖\ngithub.com/Felix201209/otto', 3, 4.8, 7, 1, size=20, color=MUTED, align=PP_ALIGN.CENTER)
cta = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(6.0), Inches(4.3), Inches(0.9))
cta.fill.solid(); cta.fill.fore_color.rgb = ACCENT
cta.line.fill.background()
tf = cta.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = 'Join Beta →'
run.font.size = Pt(22); run.font.bold = True; run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
print('  S8_CLOSE')

# Save
prs.save(OUT)
sz = os.path.getsize(OUT)
print(f'\nDONE: {OUT} ({sz//1024} KB, {len(prs.slides)} slides)')
print('All text/charts/tables/shapes are NATIVE PowerPoint objects.')
print('Open in PowerPoint → click any text → edit directly.')
print('Right-click chart → Edit Data → change values.')

# Cleanup
import shutil
shutil.rmtree(TMP, ignore_errors=True)
